import docx
from pptx import Presentation
from sentence_transformers import SentenceTransformer
import numpy as np
import torch

# Limit CPU threads to prevent system saturation on the 16-core VM
torch.set_num_threads(8)
from pypdf import PdfReader

# Lazy load model
_MODEL = None

def get_model():
    global _MODEL
    if _MODEL is None:
        try:
            # Try loading from local cache first to avoid network checks and latency
            _MODEL = SentenceTransformer("intfloat/multilingual-e5-small", local_files_only=True)
        except Exception:
            # Fallback to downloading if not cached
            _MODEL = SentenceTransformer("intfloat/multilingual-e5-small")
    return _MODEL

import logging

def extract_pdf(path):
    pages_text = []
    try:
        reader = PdfReader(str(path))
        for idx, page in enumerate(reader.pages):
            t = page.extract_text()
            # If the page has digital text of reasonable length, use it.
            if t and len(t.strip()) > 50:
                pages_text.append(t)
            else:
                # Page might be scanned/image-only. Attempt OCR.
                try:
                    import pytesseract
                    from pdf2image import convert_from_path
                    page_num = idx + 1
                    logging.info(f"Nativo arrojó texto corto o vacío en página {page_num} de {path.name}. Intentando OCR...")
                    images = convert_from_path(str(path), first_page=page_num, last_page=page_num)
                    if images:
                        ocr_text = pytesseract.image_to_string(images[0], lang='spa+eng')
                        if ocr_text and len(ocr_text.strip()) > 0:
                            logging.info(f"OCR exitoso en página {page_num}: {len(ocr_text)} caracteres reconocidos.")
                            pages_text.append(ocr_text)
                        else:
                            logging.warning(f"OCR no detectó texto en página {page_num}.")
                    else:
                        if t: pages_text.append(t)
                except Exception as e:
                    logging.error(f"Fallo de OCR en página {idx+1} de {path.name}: {e}")
                    if t:
                        pages_text.append(t)
    except Exception as e:
        logging.error(f"Error al procesar PDF {path.name}: {e}")
    return "\n".join(pages_text)

def extract_docx(path):
    paras = []
    try:
        doc = docx.Document(str(path))
        for para in doc.paragraphs:
            if para.text:
                paras.append(para.text)
    except Exception:
        pass
    return "\n".join(paras)

def extract_pptx(path):
    lines = []
    try:
        prs = Presentation(str(path))
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    lines.append(shape.text)
    except Exception:
        pass
    return "\n".join(lines)

def embed_text_chunks(chunks, prefix=""):
    # E5 models work best with "query: " or "passage: " prefixes
    if prefix:
        chunks = [f"{prefix}{c}" for c in chunks]
        
    model = get_model()
    embeddings = model.encode(chunks, convert_to_numpy=True)
    
    # L2 Normalization for Cosine Similarity
    norms = np.linalg.norm(embeddings, axis=1, keepdims=True)
    return embeddings / (norms + 1e-10)



def get_instrument_list():
    """
    Loads the instruments from the instrumentos.json file.
    Returns a list of dicts: [{'name': '...', 'definition': '...'}]
    """
    import json
    import os
    from pathlib import Path
    
    # Resolve path relative to this file
    base_dir = os.path.dirname(__file__)
    json_path = Path(base_dir) / "documentos_maestros" / "instrumentos.json"
    
    if not json_path.exists():
        return []
        
    try:
        content = json_path.read_text(encoding="utf-8")
        return json.loads(content)
    except Exception:
        return []
